from pptx import Presentation

# Create a new presentation for the customized proposal for atl Capital
prs_atl_capital = Presentation()

# Slide 1 - Introduction to Kit Consultivo
slide_layout = prs_atl_capital.slide_layouts[1]
slide = prs_atl_capital.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Presentación del Kit Consultivo para atl Capital"
content.text = """Introducción al programa Kit Consultivo, su propósito y cómo puede beneficiar a atl Capital en 
la digitalización y optimización de sus servicios financieros."""

# Slide 2 - Objetivos del Kit Consultivo en atl Capital
slide_layout = prs_atl_capital.slide_layouts[1]
slide = prs_atl_capital.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Objetivos del Kit Consultivo en atl Capital"
content.text = """- Apoyar a atl Capital en la integración de inteligencia artificial en procesos clave.
- Optimizar la eficiencia operativa en la gestión de patrimonios mediante la digitalización.
- Fortalecer el análisis y la personalización de la oferta de productos financieros a través de herramientas avanzadas."""

# Slide 3 - Casos de Uso Potenciales
slide_layout = prs_atl_capital.slide_layouts[1]
slide = prs_atl_capital.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Casos de Uso Potenciales"
content.text = """1. Inteligencia Artificial para la Personalización de Carteras
   - IA para analizar datos y crear carteras adaptadas al perfil de riesgo y objetivos.
2. Análisis Predictivo para la Gestión de Riesgos
   - IA para predecir movimientos del mercado y ajustar carteras.
3. Digitalización del Proceso de Monitorización de Inversiones
   - Automatización para ofrecer reportes en tiempo real con análisis detallado."""

# Slide 4 - Beneficios de la Implementación
slide_layout = prs_atl_capital.slide_layouts[1]
slide = prs_atl_capital.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Beneficios de la Implementación para atl Capital"
content.text = """- Optimización del tiempo y recursos mediante la automatización de procesos.
- Mejora en la experiencia del cliente con informes personalizados y monitoreo en tiempo real.
- Fortalecimiento de la competitividad en el mercado al integrar IA y digitalización en la gestión patrimonial."""

# Slide 5 - Ejemplo de Proceso con IA en Gestión Patrimonial
slide_layout = prs_atl_capital.slide_layouts[1]
slide = prs_atl_capital.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Ejemplo de Proceso con IA en Gestión Patrimonial"
content.text = """Demostración de un flujo de trabajo de IA para mejorar la asignación de activos, monitoreo y 
ajustes en las carteras de inversión de clientes de atl Capital."""

# Slide 6 - Financiación y Coste de Implementación
slide_layout = prs_atl_capital.slide_layouts[1]
slide = prs_atl_capital.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Financiación y Coste de Implementación"
content.text = """Información sobre los niveles de financiación disponibles en el Kit Consultivo, los requisitos 
para acceder a ellos y cómo atl Capital puede beneficiarse para reducir los costes de implementación."""

# Save the presentation
atl_capital_pptx_path = "Propuesta_Kit_Consultivo_atl_Capital.pptx"
prs_atl_capital.save(atl_capital_pptx_path)

atl_capital_pptx_path
